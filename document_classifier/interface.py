import streamlit as st
import pandas as pd
import plotly.express as px
import PyPDF2
import docx
import io
import requests
import time
import json
import os
from openpyxl import load_workbook
from pptx import Presentation
import csv
import zipfile
import tempfile

API_URL = os.getenv("API_URL", "http://localhost:8000")

def load_translations(lang_code):
    file_path = os.path.join('locales', f'{lang_code}.json')
    with open(file_path, 'r', encoding='utf-8') as file:
        return json.load(file)
    
translations = {
    'en': load_translations('en'),
    'tr': load_translations('tr')
}

def t(key):
    return translations[st.session_state.lang_code].get(key, key)

def create_classified_zip(results, documents):
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(os.path.join(temp_dir, 'classifications.zip'), 'w') as zipf:
            for result, document in zip(results, documents):
                classification = result['Classification'].lower()
                filename = result['Filename']
                
                folder_path = os.path.join(temp_dir, classification)
                if not os.path.exists(folder_path):
                    os.makedirs(folder_path)
                
                # TODO: Fix PDF, DOCX, PPTX, XLSX, CSV, PSV file types
                file_path = os.path.join(folder_path, filename)
                with open(file_path, 'wb') as f:
                    if isinstance(document, tuple):
                        f.write(document[1].encode('utf-8'))
                    else:
                        f.write(document.getvalue())
                
                zipf.write(file_path, os.path.join(classification, filename))
        
        with open(os.path.join(temp_dir, 'classifications.zip'), 'rb') as f:
            return f.read()

def get_installed_models():
    response = requests.get(f"{API_URL}/tags")
    if response.status_code == 200:
        models = response.json()["models"]
        return [model["name"].split(":")[0] for model in models]
    else:
        st.error(t("failed_to_fetch"))
        return []
    
def classify_document(text, model):
    lang = st.session_state.lang_code
    try:
        response = requests.post(f"{API_URL}/classify", json={"text": text, "model": model, "lang": lang})
        response.raise_for_status()
        return response.json()["classification"]
    except requests.RequestException as e:
        st.error(t("error_classifying"))
        return "Classification Error"
    
def explain_classification(text, model, classification):
    lang = st.session_state.lang_code
    try:
        response = requests.post(f"{API_URL}/explain", json={"text": text, "model": model, "classification": classification, "lang": lang})
        response.raise_for_status()
        return response.json()["explanation"]
    except requests.RequestException as e:
        st.error(f"Error explaining classification: {str(e)}")
        return "Explanation Error"

def extract_text_from_document(file):
    if isinstance(file, str):
        file_extension = os.path.splitext(file)[1].lower()
        with open(file, 'rb') as f:
            file_content = f.read()
    else:
        file_extension = os.path.splitext(file.name)[1].lower()
        file_content = file.getvalue()

    if file_extension == '.txt':
        return file_content.decode('utf-8')
    elif file_extension == '.pdf':
        return extract_text_from_pdf(file_content)
    elif file_extension in ['.docx', '.doc']:
        return extract_text_from_docx(file_content)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_content)
    elif file_extension in ['.xlsx', '.xls']:
        return extract_text_from_excel(file_content)
    elif file_extension in ['.csv', '.psv']:
        return extract_text_from_csv(file_content, delimiter=',' if file_extension == '.csv' else '|')
    else:
        st.error(f"Unsupported file type: {file_extension}")
        return ""


def extract_text_from_pdf(file_content):
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        return "\n".join([page.extract_text() for page in pdf_reader.pages])
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")
        return ""
    

def extract_text_from_docx(file_content):
    try:
        doc = docx.Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        st.error(f"Error reading DOCX: {str(e)}")
        return ""
    

def extract_text_from_pptx(file_content):
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        st.error(f"Error reading PPTX: {str(e)}")
        return ""


def extract_text_from_excel(file_content):
    try:
        wb = load_workbook(io.BytesIO(file_content), read_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join(str(cell) for cell in row if cell is not None))
        return "\n".join(text)
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return ""


def extract_text_from_csv(file_content, delimiter=','):
    try:
        text = []
        csv_data = csv.reader(io.StringIO(file_content.decode('utf-8')), delimiter=delimiter)
        for row in csv_data:
            text.append("\t".join(row))
        return "\n".join(text)
    except Exception as e:
        st.error(f"Error reading CSV/PSV file: {str(e)}")
        return ""

    
def scan_directory(directory, model):
    results = []
    documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.txt', '.pdf', '.docx', '.doc', 'pptx', 'xlsx', 'xls', 'csv', 'psv')):
                file_path = os.path.join(root, file)
                text = extract_text_from_document(file_path)
                documents.append((file, text))
                if text:
                    classification = classify_document(text, model)
                    results.append({"Filename": file, "Classification": classification})
    return results, documents


def generate_training_dataset(results, uploaded_files, scan=False):
    training_data = []

    if scan:
        for result, (file, text) in zip(results["Classification"], uploaded_files):
            training_data.append({
                "Content": text,
                "Classification": result
            })
        return pd.DataFrame(training_data)

    for result, file in zip(results["Classification"], uploaded_files):
        text = extract_text_from_document(file)
        training_data.append({
            "Content": text,
            "Classification": result
        })
    return pd.DataFrame(training_data)


st.set_page_config(page_title="Document Classification System", layout="wide")
lang = st.selectbox("Select Language / Dil Seçin", ["English", "Türkçe"], index=0, key="lang")
st.session_state.lang_code = 'en' if lang == "English" else 'tr'

if "showresults" not in st.session_state:
    st.session_state.showresults = False

if "results" not in st.session_state:
    st.session_state.results = []

if "showscan" not in st.session_state:
    st.session_state.showscan = False

if "documents" not in st.session_state:
    st.session_state.documents = []

st.title(t("title"))

installed_models = get_installed_models()
local_models = ["gemma2", "phi3", "llama3", "mistral", "llama3.1", "mistral-nemo", "mertergun/phi3_finetuned", "gemma2", "qwen2"]
all_models = installed_models + [m for m in local_models if m not in installed_models] + ["gpt-3.5-turbo-0125", "gpt-3.5-turbo-instruct", "gpt-4-turbo", "gpt-4o", "gemini-1.5-pro", "gemini-1.5-flash", "gemini-1.0-pro", "gpt-4o-mini"]

model = st.selectbox(t("model_select"), all_models)
if model not in installed_models and model in local_models:
            if st.button(t("pull_model")):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                def update_progress(completed, total):
                    progress = (completed / total) * 100 if total > 0 else 0
                    progress = min(progress, 100)  # Ensure progress doesn't exceed 100%
                    progress_bar.progress(int(progress))
                    status_text.text(f"Pulling model... {progress:.2f}% ({completed}/{total} bytes)")

                response = requests.post(f"{API_URL}/pull_model", json={"model": model}, stream=True)
                if response.status_code == 200:
                    total_size = None
                    completed_size = 0
                    for line in response.iter_lines():
                        if line:
                            try:
                                data = json.loads(line)
                                if 'total' in data:
                                    total_size = int(data['total'])
                                if 'completed' in data:
                                    completed_size = int(data['completed'])
                                    if total_size:
                                        update_progress(completed_size, total_size)
                            except json.JSONDecodeError:
                                status_text.text(line.decode('utf-8'))
                    
                    if total_size and completed_size == total_size:
                        st.success(f"Model '{model}' has been successfully pulled.")
                    elif total_size:
                        st.warning(f"Model pull may not have completed. Please check the model status.")
                    else:
                        st.info(f"Model '{model}' pull process completed, but size information was not available.")
                    
                    time.sleep(1)  
                    st.rerun()
                else:
                    st.error("Failed to pull the selected model.")

# Set a radio button to select the select upload files or scan a directory
option = st.radio(t("option_radio"), [t("upload_files"), t("scan_directory")])

if option == t("scan_directory"):
    st.write(t("please_select"))
    directory = st.text_input(t("directory_path"), value=".")
    if st.button(t("scan_directory"), key="scan_directory") or st.session_state.showscan:
        if not st.session_state.showscan or st.session_state.scan_directory:
            st.session_state.showscan = True
            with st.spinner(f"{t('scanning_directory')} {directory}..."):
                st.session_state.results, st.session_state.documents = scan_directory(directory, model)
        
        if st.session_state.results:
            results_df = pd.DataFrame(st.session_state.results)
            st.subheader(t("classification_results"))
            # Make the dataframe editable, so user change the classification if needed using selectbox in the Classification column
            edited_df = st.data_editor(results_df, use_container_width=True, column_config={
                "Classification": st.column_config.SelectboxColumn(
                    "Classification",
                    width="medium",
                    options=[t("ts"), t("s"), t("c"), t("r"), t("u")],
                    required=True,
                )
            }, num_rows="fixed")
            
            # Visualization
            fig = px.pie(edited_df, names="Classification", title=t('distribution'))
            st.plotly_chart(fig)
        
            # Download results
            csv = results_df.to_csv(index=False)
            st.download_button(
                label=t("download_results"),
                data=csv,
                file_name="classification_results.csv",
                mime="text/csv",
            )

            if st.button(t("generate_dataset")):
                with st.spinner(t("generating_dataset")):
                    training_data = generate_training_dataset(edited_df, st.session_state.documents, scan=True)
                    csv = training_data.to_csv(index=False)
                    st.download_button(
                        label="Download Train Dataset",
                        data=csv,
                        file_name="train_dataset.csv",
                        mime="text/csv",
                    )

            if st.button(t("download_classified_zip")):
                zip_file = create_classified_zip(st.session_state.results, st.session_state.documents)
                st.download_button(
                    label=t("download_classified_documents"),
                    data=zip_file,
                    file_name="classifications.zip",
                    mime="application/zip",
                )
        else:
            st.warning(t("no_valid_documents"))

else:
    uploaded_files = st.file_uploader(t("upload_document"), accept_multiple_files=True, type=['txt', 'pdf', 'docx', 'doc', 'pptx', 'xlsx', 'xls', 'csv', 'psv'])

    if st.button(t("upload_files"), key="upload") or st.session_state.showresults:
        if not st.session_state.showresults or st.session_state.upload:
            st.session_state.showresults = True
            st.session_state.results = []
            for file in uploaded_files:
                with st.spinner(f'Processing {file.name}...'):
                    text = extract_text_from_document(file)
                    if text:
                        classification = classify_document(text, model)
                        st.session_state.results.append({"Filename": file.name, "Classification": classification})
        
        if st.session_state.results:
            results_df = pd.DataFrame(st.session_state.results)
            
            if len(st.session_state.results) == 1:
                st.success(t(f"completed"))
                st.markdown(f"<h2 style='text-align: center; color: #1E90FF;'>{t('classification_results')}</h2>", unsafe_allow_html=True)
                st.markdown(f"<h3 style='text-align: center;'>{t('file')}: {st.session_state.results[0]['Filename']}</h3>", unsafe_allow_html=True)
                
                classification = st.session_state.results[0]['Classification']
                color = "#000000" 
                if classification.lower() == t("ts"):
                    color = "#FF0000"  
                elif classification.lower() == t("s"):
                    color = "#FFA500"  
                elif classification.lower() == t("c"):
                    color = "#800080" 
                elif classification.lower() == t("r"):
                    color = "#0000FF" 
                elif classification.lower() == t("u"):
                    color = "#008000"  
                else: # Means error
                    color = "#000000"
                
                st.markdown(f"<h1 style='text-align: center; color: {color};'>{classification.upper()}</h1>", unsafe_allow_html=True)

                if uploaded_files:
                    text = extract_text_from_document(uploaded_files[0])

                    if st.button(t("explain")):
                        explanation = explain_classification(text, model, classification)
                        st.write(explanation)

            else:
                st.subheader(t("classification_results"))
                # Make the dataframe editable, so user change the classification if needed using selectbox in the Classification column
                edited_df = st.data_editor(results_df, use_container_width=True, column_config={
                    "Classification": st.column_config.SelectboxColumn(
                        "Classification",
                        width="medium",
                        options=[t("ts"), t("s"), t("c"), t("r"), t("u")],
                        required=True,
                    )
                }, num_rows="fixed")
                
                # Visualization
                fig = px.pie(edited_df, names="Classification", title=t('distribution'))
                st.plotly_chart(fig)
            
                # Download results
                csv = results_df.to_csv(index=False)
                st.download_button(
                    label=t("download_results"),
                    data=csv,
                    file_name="classification_results.csv",
                    mime="text/csv",
                )

                if st.button(t("generate_dataset")):
                    with st.spinner(t("generating_dataset")):
                        training_data = generate_training_dataset(edited_df, uploaded_files)
                        csv = training_data.to_csv(index=False)
                        st.download_button(
                            label="Download Train Dataset",
                            data=csv,
                            file_name="train_dataset.csv",
                            mime="text/csv",
                        )

                if st.button(t("download_classified_zip")):
                    zip_file = create_classified_zip(st.session_state.results, uploaded_files)
                    st.download_button(
                        label=t("download_classified_documents"),
                        data=zip_file,
                        file_name="classifications.zip",
                        mime="application/zip",
                    )

        else:
            st.warning(t("no_valid_documents"))

st.sidebar.title(t("about"))
st.sidebar.info(
    f"{t('about1')}\n"
    f"{t('about2')}\n"
    f"{t('about3')}\n"
    f"{t('about4')}"
)

st.sidebar.title(t("instructions"))
st.sidebar.write(
    f"{t('instructions1')}\n"
    f"{t('instructions2')}\n"
    f"{t('instructions3')}\n"
    f"{t('instructions4')}\n"
    f"{t('instructions5')}\n"
    f"{t('instructions6')}"
)