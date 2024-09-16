import os
import io
import csv
import zipfile
import tempfile
import PyPDF2
import docx
from openpyxl import load_workbook
from pptx import Presentation
import xlrd 
import pandas as pd

from document_classifier.utils.interface_utils import classify_document

def create_classified_zip(results, documents):
    """
    Creates a zip file containing documents organized into folders based on 
    their classification.
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(
            os.path.join(temp_dir, 'classifications.zip'), 'w'
        ) as zipf:
            for result, document in zip(results, documents):
                classification = result['Classification'].lower()
                filename = result['Filename']
                
                folder_path = classification

                if isinstance(document, tuple):
                    file_name, file_content = document
                    if isinstance(file_content, str):
                        file_content = file_content.encode('utf-8')
                    zipf.writestr(
                        os.path.join(folder_path, file_name), file_content
                    )
                elif hasattr(document, 'read'):
                    document.seek(0)
                    zipf.writestr(
                        os.path.join(folder_path, filename), document.read()
                    )
                else:
                    if isinstance(document, str):
                        document = document.encode('utf-8')
                    zipf.writestr(
                        os.path.join(folder_path, filename), document
                    )
        
        with open(os.path.join(temp_dir, 'classifications.zip'), 'rb') as f:
            return f.read()

def extract_text_from_document(file, file_extension=None):
    """Extracts text from various document types."""
    if isinstance(file, bytes):
        file_content = file
        if not file_extension:
            file_extension = guess_file_extension(file_content)
    elif isinstance(file, str):
        file_extension = os.path.splitext(file)[1].lower()
        with open(file, 'rb') as f:
            file_content = f.read()
    else:
        file_extension = os.path.splitext(file.name)[1].lower()
        file_content = file.getvalue()

    if file_extension == '.txt':
        return file_content.decode('utf-8', errors='ignore')
    elif file_extension == '.pdf':
        return extract_text_from_pdf(file_content)
    elif file_extension in ['.docx', '.doc']:
        return extract_text_from_docx(file_content)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_content)
    elif file_extension == '.xlsx':
        return extract_text_from_excel(file_content)
    elif file_extension == '.xls':
        return extract_text_from_xls(file_content)
    elif file_extension in ['.csv', '.psv']:
        return extract_text_from_csv(
            file_content, delimiter=',' if file_extension == '.csv' else '|'
        )
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

def guess_file_extension(file_content):
    """Guesses the file extension based on content."""
    if file_content.startswith(b'%PDF'):
        return '.pdf'
    elif file_content.startswith(b'PK\x03\x04'):
        return '.docx' 
    elif file_content.startswith(b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'):
        return '.doc' 
    else:
        return '.txt'

def extract_text_from_pdf(file_content):
    """Extracts text from PDF files."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        return "\n".join([page.extract_text() for page in pdf_reader.pages])
    except Exception as e:
        raise ValueError(f"Error reading PDF: {str(e)}") 

def extract_text_from_docx(file_content):
    """Extracts text from DOCX files."""
    try:
        doc = docx.Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ValueError(f"Error reading DOCX: {str(e)}") 

def extract_text_from_pptx(file_content):
    """Extracts text from PPTX files."""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error reading PPTX: {str(e)}")

def extract_text_from_excel(file_content):
    """Extracts text from XLSX files."""
    try:
        wb = load_workbook(io.BytesIO(file_content), read_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text.append(
                    "\t".join(str(cell) for cell in row if cell is not None)
                )
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error reading Excel file: {str(e)}") 

def extract_text_from_xls(file_content):
    """Extracts text from XLS files."""
    try:
        wb = xlrd.open_workbook(file_contents=file_content)
        text = []
        for sheet in wb.sheet_names():
            ws = wb.sheet_by_name(sheet)
            for row in range(ws.nrows):
                text.append(
                    "\t".join(
                        str(cell) for cell in ws.row_values(row) if cell is not None
                    )
                )
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error reading XLS file: {str(e)}")


def extract_text_from_csv(file_content, delimiter=','):
    """Extracts text from CSV/PSV files."""
    try:
        text = []
        csv_data = csv.reader(
            io.StringIO(file_content.decode('utf-8')), delimiter=delimiter
        )
        for row in csv_data:
            text.append("\t".join(row))
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error reading CSV/PSV file: {str(e)}")

def scan_directory(directory, model):
    results = []
    documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.txt', '.pdf', '.docx', '.doc', 'pptx', 'xlsx', 'xls', 'csv', 'psv')):
                file_extension = os.path.splitext(file)[1].lower()
                file_path = os.path.join(root, file)
                with open(file_path, 'rb') as f:
                    file_content = f.read()
                text = extract_text_from_document(file_content, file_extension)
                documents.append((file, file_content))
                if text:
                    classification = classify_document(text, model)
                    results.append({"Filename": file, "Classification": classification})
    return results, documents

def generate_training_dataset(results, uploaded_files, scan=False):
    training_data = []

    if scan:
        for result, (file, text) in zip(results["Classification"], uploaded_files):
            training_data.append({
                "Content": text,
                "Classification": result
            })
        return pd.DataFrame(training_data)

    for result, file in zip(results["Classification"], uploaded_files):
        text = extract_text_from_document(file)
        training_data.append({
            "Content": text,
            "Classification": result
        })
    return pd.DataFrame(training_data)

